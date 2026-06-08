from __future__ import annotations

import posixpath
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Callable, Iterable


Normalize = Callable[[object], str]


def _dedupe_append(blocks: list[dict], seen: set[tuple[int, str, str]], block: dict, normalize: Normalize) -> None:
    text = normalize(block.get("text", ""))
    if not text:
        return
    key = (int(block.get("slide", 0) or 0), str(block.get("type", "")), re.sub(r"\s+", " ", text))
    if key in seen:
        return
    seen.add(key)
    block["text"] = text
    blocks.append(block)


def _iter_shape_text_blocks(shape: object, slide_index: int, title_shape: object | None, normalize: Normalize) -> Iterable[dict]:
    if hasattr(shape, "shapes"):
        for child in getattr(shape, "shapes", []):
            yield from _iter_shape_text_blocks(child, slide_index, title_shape, normalize)
        return

    if getattr(shape, "has_table", False):
        table = getattr(shape, "table", None)
        if table is not None:
            for row_index, row in enumerate(table.rows, start=1):
                cells = [normalize(cell.text) for cell in row.cells]
                cells = [cell for cell in cells if cell]
                if cells:
                    yield {
                        "type": "table_row",
                        "style": "PPTXTable",
                        "slide": slide_index,
                        "row": row_index,
                        "text": " | ".join(cells),
                    }

    if not getattr(shape, "has_text_frame", False):
        return

    paragraphs: list[str] = []
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is not None:
        for paragraph in text_frame.paragraphs:
            text = normalize("".join(run.text for run in paragraph.runs) or paragraph.text)
            if text:
                paragraphs.append(text)

    if not paragraphs:
        return

    is_title = title_shape is not None and getattr(shape, "element", None) is getattr(title_shape, "element", None)
    block_type = "slide_title" if is_title else "paragraph"
    style = getattr(shape, "name", "") or "PPTXText"
    yield {
        "type": block_type,
        "style": style,
        "slide": slide_index,
        "text": "\n".join(paragraphs),
    }


def _extract_with_python_pptx(path: Path, normalize: Normalize) -> list[dict]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError(f"PPTX parser import failed (python-pptx required): {exc}") from exc

    prs = Presentation(str(path))
    blocks: list[dict] = []
    seen: set[tuple[int, str, str]] = set()
    for slide_index, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            for block in _iter_shape_text_blocks(shape, slide_index, title_shape, normalize):
                _dedupe_append(blocks, seen, block, normalize)
    return blocks


def _join_xml_text(root: ET.Element, normalize: Normalize) -> str:
    parts = [node.text or "" for node in root.findall(".//{*}t")]
    return normalize(" ".join(part for part in parts if part))


def _paragraph_xml_texts(root: ET.Element, normalize: Normalize) -> list[str]:
    texts: list[str] = []
    for paragraph in root.findall(".//{*}p"):
        parts = [node.text or "" for node in paragraph.findall(".//{*}t")]
        text = normalize(" ".join(part for part in parts if part))
        if text:
            texts.append(text)
    return texts


def _resolve_pptx_target(source_part: str, target: str) -> str:
    base_dir = posixpath.dirname(source_part)
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(base_dir, target))


def _load_slide_note_paths(zf: zipfile.ZipFile, slide_path: str) -> list[str]:
    rels_path = f"{posixpath.dirname(slide_path)}/_rels/{posixpath.basename(slide_path)}.rels"
    try:
        root = ET.fromstring(zf.read(rels_path))
    except KeyError:
        return []

    note_paths: list[str] = []
    for rel in root.findall("{*}Relationship"):
        rel_type = rel.attrib.get("Type", "")
        target = rel.attrib.get("Target", "")
        if target and rel_type.endswith("/notesSlide"):
            note_paths.append(_resolve_pptx_target(slide_path, target))
    return note_paths


def _extract_notes_blocks(path: Path, normalize: Normalize) -> list[dict]:
    blocks: list[dict] = []
    with zipfile.ZipFile(path) as zf:
        slide_paths = sorted(
            (name for name in zf.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
            key=lambda name: int(re.search(r"slide(\d+)\.xml$", name).group(1)),  # type: ignore[union-attr]
        )
        for slide_index, slide_path in enumerate(slide_paths, start=1):
            for note_path in _load_slide_note_paths(zf, slide_path):
                try:
                    root = ET.fromstring(zf.read(note_path))
                except KeyError:
                    continue
                text = _join_xml_text(root, normalize)
                if text and text.lower() != "click to add notes":
                    blocks.append(
                        {
                            "type": "speaker_notes",
                            "style": "PPTXNotes",
                            "slide": slide_index,
                            "text": text,
                        }
                    )
    return blocks


def _extract_with_raw_ooxml(path: Path, normalize: Normalize) -> list[dict]:
    blocks: list[dict] = []
    with zipfile.ZipFile(path) as zf:
        slide_paths = sorted(
            (name for name in zf.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
            key=lambda name: int(re.search(r"slide(\d+)\.xml$", name).group(1)),  # type: ignore[union-attr]
        )
        for slide_index, slide_path in enumerate(slide_paths, start=1):
            root = ET.fromstring(zf.read(slide_path))
            for para_index, text in enumerate(_paragraph_xml_texts(root, normalize), start=1):
                blocks.append(
                    {
                        "type": "paragraph",
                        "style": "PPTXRawXML",
                        "slide": slide_index,
                        "paragraph": para_index,
                        "text": text,
                    }
                )
    blocks.extend(_extract_notes_blocks(path, normalize))
    return blocks


def extract_pptx_blocks(path: Path, normalize: Normalize) -> list[dict]:
    """Extract searchable text blocks from a PowerPoint deck.

    python-pptx is used first for shape/table semantics. If it cannot parse a
    deck, the function falls back to direct OOXML extraction.
    """
    blocks: list[dict] = []
    seen: set[tuple[int, str, str]] = set()

    try:
        blocks.extend(_extract_with_python_pptx(path, normalize))
    except Exception:
        blocks = _extract_with_raw_ooxml(path, normalize)
    else:
        for block in blocks:
            text = normalize(block.get("text", ""))
            if text:
                seen.add((int(block.get("slide", 0) or 0), str(block.get("type", "")), re.sub(r"\s+", " ", text)))
        for block in _extract_notes_blocks(path, normalize):
            _dedupe_append(blocks, seen, block, normalize)

    return [block for block in blocks if normalize(block.get("text", ""))]


__all__ = ["extract_pptx_blocks"]
